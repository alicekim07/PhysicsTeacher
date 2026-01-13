import subprocess
from pathlib import Path
from pptx import Presentation

def pptx_to_pdf(pptx_path: str, pdf_dir: str):
    pdf_dir = Path(pdf_dir)
    pdf_dir.mkdir(parents=True, exist_ok=True)

    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        pptx_path,
        "--outdir", str(pdf_dir)
    ], check=True)

def extract_notes_from_pptx(pptx_path:str) -> str:
    prs = Presentation(pptx_path)
    notes = []

    for idx, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text
            if text.strip():
                notes.append(f"[Slide {idx+1}]\n{text.strip()}")

    return "\n\n".join(notes)